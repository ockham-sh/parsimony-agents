"""Tests for workspace document helpers."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pypdf import PdfWriter

from parsimony_agents.execution.documents import read_excel, read_pdf_text, read_pptx_text


def test_read_pdf_text_empty_pages(tmp_path: Path) -> None:
    p = tmp_path / "t.pdf"
    w = PdfWriter()
    w.add_blank_page(width=72, height=72)
    with p.open("wb") as f:
        w.write(f)
    assert read_pdf_text(str(p)) == ""


def test_read_excel_roundtrip(tmp_path: Path) -> None:
    p = tmp_path / "s.xlsx"
    pd.DataFrame({"a": [1, 2]}).to_excel(p, index=False, engine="openpyxl")
    df = read_excel(str(p))
    assert list(df.columns) == ["a"]
    assert len(df) == 2


def test_read_pptx_text_roundtrip(tmp_path: Path) -> None:
    from pptx import Presentation

    p = tmp_path / "d.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(0, 0, 100, 50)
    box.text_frame.text = "Hello deck"
    prs.save(str(p))

    slides = read_pptx_text(str(p))
    assert len(slides) >= 1
    assert "Hello deck" in slides[0]["text"]
