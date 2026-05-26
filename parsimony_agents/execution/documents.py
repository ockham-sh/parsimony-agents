"""Workspace document helpers (PDF, Excel, PowerPoint) for notebook kernels.

Imports run at call time so ``parsimony-agents`` can be installed without
the optional ``[documents]`` extra; the executor injects these callables
into ``locals`` (users do not ``import`` third-party packages directly).
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

import pandas as pd


def _require_pypdf():
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise RuntimeError(
            "PDF support requires the optional documents stack: install parsimony-agents with the "
            "``documents`` extra (included in ``[all]``)."
        ) from e
    return PdfReader


def read_pdf_text(path: str, *, max_pages: int | None = None) -> str:
    """Extract plain text from a PDF under the workspace (relative *path*)."""
    PdfReader = _require_pypdf()
    p = Path(path)
    with p.open("rb") as f:
        data = f.read()
    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    n = len(reader.pages)
    end = n if max_pages is None else min(n, max_pages)
    for i in range(end):
        page = reader.pages[i]
        t = page.extract_text()
        if t and t.strip():
            parts.append(t)
    return "\n\n".join(parts)


def _require_openpyxl_for_excel() -> None:
    try:
        import openpyxl  # noqa: F401
    except ImportError as e:
        raise RuntimeError(
            "Excel support requires openpyxl: install parsimony-agents with the ``documents`` extra."
        ) from e


def read_excel(
    path: str,
    *,
    sheet_name: int | str = 0,
    **kwargs: Any,
) -> pd.DataFrame:
    """Read an Excel file into a :class:`pandas.DataFrame` (xlsx; engine openpyxl)."""
    _require_openpyxl_for_excel()
    return pd.read_excel(path, sheet_name=sheet_name, engine="openpyxl", **kwargs)


def _require_pptx():
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError(
            "PowerPoint support requires python-pptx: install parsimony-agents with the ``documents`` extra."
        ) from e
    return Presentation


def read_pptx_text(path: str) -> list[dict[str, Any]]:
    """Return a list of per-slide text extracts from a ``.pptx`` file.

    Each item has ``index`` (0-based) and ``text`` (joined shapes text).
    """
    Presentation = _require_pptx()
    prs = Presentation(path)
    out: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                t = p.text
                if t and t.strip():
                    parts.append(t.strip())
        out.append({"index": i, "text": "\n".join(parts)})
    return out
